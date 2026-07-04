from pptx import Presentation

from ..models import ExtractionResult, Segment


def extract(path: str) -> ExtractionResult:
    prs = Presentation(path)
    segments: list[Segment] = []
    refs: dict[str, object] = {}
    order = 0

    def add_segment(paragraph, group_key: str) -> None:
        nonlocal order
        text = paragraph.text
        if not text.strip():
            return
        seg = Segment(doc_type="pptx", source_text=text, group_key=group_key, order_hint=order)
        order += 1
        segments.append(seg)
        refs[seg.id] = paragraph

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    add_segment(para, f"slide{slide_idx}")
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            add_segment(para, f"slide{slide_idx}_table")

    return ExtractionResult(doc_handle=prs, segments=segments, refs=refs)
