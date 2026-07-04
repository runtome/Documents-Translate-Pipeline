from pptx import Presentation
from pptx.util import Inches

from doctranslate.pipeline import run_pipeline


def test_pptx_roundtrip(tmp_path, fake_llm_client):
    src = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Hello world"
    body = slide.placeholders[1]
    body.text_frame.text = "Second line"
    prs.save(src)

    out_path = run_pipeline(str(src), "en", "th", fake_llm_client)

    assert out_path.exists()
    out_prs = Presentation(str(out_path))
    assert len(out_prs.slides) == 1
    out_slide = out_prs.slides[0]
    assert out_slide.shapes.title.text == "[TR]Hello world"
    assert out_slide.placeholders[1].text_frame.text == "[TR]Second line"


def test_pptx_roundtrip_preserves_table(tmp_path, fake_llm_client):
    src = tmp_path / "with_table.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    graphic_frame = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    table = graphic_frame.table
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Role"
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "Engineer"
    prs.save(src)

    out_path = run_pipeline(str(src), "en", "ja", fake_llm_client)

    out_prs = Presentation(str(out_path))
    out_table = out_prs.slides[0].shapes[0].table
    assert out_table.cell(0, 0).text == "[TR]Name"
    assert out_table.cell(1, 1).text == "[TR]Engineer"
